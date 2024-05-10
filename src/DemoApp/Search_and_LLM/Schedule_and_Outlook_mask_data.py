import speech_recognition as sr
import os
import openai
import pyttsx3
import re

import signal
signal.signal(signal.SIGINT, signal.SIG_DFL)

from dotenv import load_dotenv
# .envファイルの内容を読み込見込む
load_dotenv()
# # os.environを用いて環境変数を表示させます
# print(os.environ['OpenAI_API_KEY'])
# key = os.environ['OpenAI_API_KEY']

##############
# 音声認識関数 #
##############
def recognize_speech():

    recognizer = sr.Recognizer()    
    # Set timeout settings.
    recognizer.dynamic_energy_threshold = False

    
    with sr.Microphone() as source:
        recognizer.adjust_for_ambient_noise(source)
    
        while(True):
            print(">> Please speak now...")

            # engine.say("Please speak now")


            audio = recognizer.listen(source, timeout=1000.0)

            try:
                # Google Web Speech API を使って音声をテキストに変換
                text = recognizer.recognize_google(audio, language="ja-JP")
                print("[You]")
                print(text)
                return text
            except sr.UnknownValueError:
                print("Sorry, I could not understand what you said. Please speak again.")
                #return ""
            except sr.RequestError as e:
                print(f"Could not request results; {e}")
                #return ""


#################################
# Pyttsx3でレスポンス内容を読み上げ #
#################################
def text_to_speech(text):
    # テキストを読み上げる
    engine.say(text)
    engine.runAndWait()



def chat(conversationHistory):
    # APIリクエストを作成する
    response = openai.chat.completions.create(
        messages=conversationHistory,
        max_tokens=1024,
        n=1,
        stream=True,
        temperature=0.5,
        stop=None,
        presence_penalty=0.5,
        frequency_penalty=0.5,
        model="gpt-3.5-turbo"
    )

    # ストリーミングされたテキストを処理する
    fullResponse = ""
    RealTimeResponce = ""   
    for chunk in response:
        text = chunk.choices[0].delta.content # chunk['choices'][0]['delta'].get('content')

        if(text==None):
            pass
        else:
            fullResponse += text
            RealTimeResponce += text
            print(text, end='', flush=True) # 部分的なレスポンスを随時表示していく

            target_char = ["。", "！", "？", "\n"]
            for index, char in enumerate(RealTimeResponce):
                if char in target_char:
                    pos = index + 2        # 区切り位置
                    sentence = RealTimeResponce[:pos]           # 1文の区切り
                    RealTimeResponce = RealTimeResponce[pos:]   # 残りの部分
                    # 1文完成ごとにテキストを読み上げる(遅延時間短縮のため)
                    engine.say(sentence)
                    engine.runAndWait()
                    break
                else:
                    pass

    # APIからの完全なレスポンスを返す
    return fullResponse


##############
# メインの関数 #
##############
if __name__ == '__main__':

    ##################
    # ChatGPTの初期化 #
    ##################
    # openai.api_key=key # 自身のAPIキーを指定"
    # UserとChatGPTとの会話履歴を格納するリスト
    conversationHistory = []
    # setting = {"role": "system", "content": "句読点と読点を多く含めて応答するようにして下さい。また、1文あたりが長くならないようにして下さい。"}
    # 通常バージョン
    setting = {"role": "system", "content": "句読点と読点を多く含めて応答するようにして下さい。また、1文あたりが長くならないようにして下さい。100文字以内でお願いします。"}
    # 箇条書きバージョン
    setting = {"role": "system", "content": "句読点と読点を多く含めて箇条書きで応答するようにして下さい。また、1文あたりが長くならないようにして下さい。100文字以内でお願いします。"}

    # setting = {"role": "user", 
    #            "content": "本日の予定として以下の情報を提供します。\
    #             あとで次の予定を聞くので、そのタイミング次の予定をリマインドしてください。\
    #             時間, 予定の内容, 場所の順で記述します。\
    #             10:00, デイリーミーティング, オンライン\
    #             12:00, ランチ, 食堂\
    #             15:00, A社と会議, 会議室A\
    #             17:00, 社内メンバーと会議, 会議室B\
    #             19:00, 帰宅, None"
    #            }
    
    import win32com.client
    import datetime
    dt_now_str = datetime.datetime.now()
    dt_now_str = dt_now_str.strftime('%Y年%m月%d日 %H:%M:%S')
    print("時間:", dt_now_str)

    Outlook = win32com.client.Dispatch("Outlook.Application").GetNameSpace("MAPI")
    items = Outlook.GetDefaultFolder(9).Items
    # 定期的な予定の二番目以降の予定を検索に含める
    items.IncludeRecurrences = True
    # 開始時間でソート
    items.Sort("[Start]")
    select_items = [] # 指定した期間内の予定を入れるリスト
    dt_now = datetime.datetime.now()
    start_date = datetime.date(dt_now.year, dt_now.month, dt_now.day)
    end_date = datetime.date(dt_now.year, dt_now.month, dt_now.day + 1)
    strStart = start_date.strftime('%m/%d/%Y %H:%M %p')
    strEnd = end_date.strftime('%m/%d/%Y %H:%M %p')
    sFilter = f"[Start] >= '{strStart}' And [End] <= '{strEnd}'"
    # フィルターを適用し表示
    FilteredItems = items.Restrict(sFilter)
    for item in FilteredItems:
        if start_date <= item.start.date() <= end_date:
            select_items.append(item)
            
    print("今日の予定の件数:", len(select_items))
    # 抜き出した予定の詳細を表示

    # ダミーデータ
    i = 0
    work = ['A社と会議', 'Bさんと面談', 'Cの開発定例', 'D社との商談', 'Eさんの資料のレビュー', 'Fチーム定例会', 'Gチーム戦略会議', 'H部定例', 'I課定例', 'J退勤', 'Kジム']
    room = ['会議室A', '会議室B', 'TeamsC', '会議室D', 'TeamsE', 'TeamsF', '会議室G', '会議室H', 'ITeams', '移動中J', 'ジムK']
    for select_item in select_items:
        # print("件名：", select_item.subject)
        print(f"件名：予定 {work[i]}")  # 社外秘情報は伏せる
        print("場所：", select_item.location)
        # print(f"場所：会議室 {room[i]}") # 社外秘情報は伏せるd
        print("開始時刻：", str(select_item.Start.Format("%Y/%m/%d %H:%M")))
        print("終了時刻：", str(select_item.End.Format("%Y/%m/%d %H:%M")))
        # print("本文：", select_item.body)
        
        print("----")
        i += 1 # ダミーデータ用

    
    # ダミーデータ
    i = 0
    # text = "今日は" +  dt_now_str + "です。"\
    #          + "本日" + dt_now_str + "の予定として以下の情報を提供します。\
    #         あとで次の予定を聞くので、そのタイミングで次の予定をリマインドしてください。\
    #         " # 予定, 場所, 時間の順で記述します。" #  + OutlookSchedule
    text = "今日は" +  dt_now_str + "です。"\
        + "本日" + dt_now_str + "の予定として以下の情報を提供します。\
        あとで予定を聞くので、そのタイミングでリマインドしてください。"
    
    
    # for item in FilteredItems:
    for select_item in select_items:
        # text += "件名：" + select_item.subject
        text += "件名：" + work[i] # 社外秘情報は伏せる
        text += "場所：" + select_item.location
        # text += "場所：" + "会議室" + room[i] # 社外秘情報は伏せる
        text += "開始時刻：" + str(select_item.Start.Format("%Y/%m/%d %H:%M"))

        text += "----"
        i += 1 # ダミーデータ用
        
    user_action = {"role": "user", "content": text}
    conversationHistory.append(user_action)





    
    # 2024/04/02 追加(試しに追加)
    # 最初に設定を入力
    # conversationHistory.append(setting)
    # 2024/04/02 追加(試しに追加)

    ##################
    # Pyttsx3を初期化 #
    ##################
    engine = pyttsx3.init()
    # 読み上げの速度を設定する
    rate = engine.getProperty('rate')
    engine.setProperty('rate', rate)
    # #rate デフォルト値は200
    # rate = engine.getProperty('rate')
    # engine.setProperty('rate',200)

    #volume デフォルト値は1.0、設定は0.0~1.0
    volume = engine.getProperty('volume')
    engine.setProperty('volume',1.0)


    # Kyokoさんに喋ってもらう(日本語)
    engine.setProperty('voice', "com.apple.ttsbundle.Kyoko-premium")

    # Ctrl-Cで中断されるまでChatGPT音声アシスタントを起動
    while True:
        
        # engine.say("Detect Combo !!")
        engine.say("Detect Gesture Combination !!")
        # engine.say("Please speak now")
        # engine.say("Please speak to me in three seconds.")
        engine.say("Please speak in three seconds.")
        engine.runAndWait()
        
        # 音声認識関数の呼び出し
        # text_speach = recognize_speech()
        
        # 定型文(プリセット)にする場合
        
        time = datetime.datetime.now()
        time = time.strftime('%H:%M:%S')
        pre_Info = "現在時刻は" + time + "です。\n" # dt_now + "です。"
        
        # LLMに入力
        text = pre_Info + "この後の予定は何ですか？何分後にどこに向かえばいいですか？" # 今日の15:00の予定は何ですか？どこに向かえばいいですか？
        # text = pre_Info + "今日の午後の予定は何ですか？\n時間と場所も教えて。"
        # text = pre_Info + text_speach

        # 2024/05/09 追加
        text += "直近に必要な情報のみ簡潔に教えてください。"

        if text:
            print(" >> Waiting for response from ChatGPT...")
            # ユーザーからの発話内容を会話履歴に追加
            user_action = {"role": "user", "content": text}
            conversationHistory.append(user_action)
            
            print("[ChatGPT]") #応答内容をコンソール出力
            res = chat(conversationHistory)
            
            # ChatGPTからの応答内容を会話履歴に追加
            chatGPT_responce = {"role": "assistant", "content": res}
            conversationHistory.append(chatGPT_responce) 
            
            # 対話の履歴を表示
            # print(conversationHistory)

            # 1回で終了する場合 ... 今は一回で終了
            break
    
    engine.say("Guidance End")
    engine.runAndWait()




    # # 2024/04/26 追加
    # # 会話内容のログを取ったり、調査レポートをまとめたテキストをもとにパワポ作成
    # import subprocess
    # "パワポ作成"
    # # subprocess.run([f'C:/Users/0107409377/.pyenv/pyenv-win/versions/3.12.0/python.exe', 'C:/Users/0107409377/Desktop/code/AtCoder/src/DemoApp/Search_and_LLM/Outlook_Schedule/PowerPoint.py', '{conversationHistory}'])
    # subprocess.run([f'C:/Users/0107409377/.pyenv/pyenv-win/versions/3.12.0/python.exe', 'C:/Users/0107409377/Desktop/code/AtCoder/src/DemoApp/Search_and_LLM/Outlook_Schedule/PowerPoint.py', {chatGPT_responce}])
    # "メールに添付して送信"
    # subprocess.run([f'C:/Users/0107409377/.pyenv/pyenv-win/versions/3.12.0/python.exe', 'C:/Users/0107409377/Desktop/code/AtCoder/src/DemoApp/Search_and_LLM/Outlook_Schedule/Outlook_Schedule/Mail.py'])
    
    "引数での渡し方がわからなかったので、代替のやり方（直接ここにパワポ作成の処理を書く）"
    from pptx import Presentation
    prs=Presentation()
    from pptx.chart.data import CategoryChartData
    c_data=CategoryChartData()
    c_data.categories=['A社','B社','C社','D社','E社','F社','G社','H社','I社', 'その他']
    c_data.add_series('xxxx（円）',(10,15,8, 5, 17, 12, 2, 8, 15, 7))
    sld0=prs.slides.add_slide(prs.slide_layouts[5]) # title only
    sld0.shapes[0].text='xxxxに関する市場調査' # 見出し
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Cm
    sld0.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        # Cm(1),Cm(3),Cm(20),Cm(15),c_data)
        Cm(1),Cm(3),Cm(15),Cm(15),c_data)
    txBox = sld0.shapes.add_textbox(Cm(15),Cm(5),Cm(3),Cm(3))
    tf = txBox.text_frame		# TextFrameオブジェクトの設定
    tf.text = "This is text inside a textbox"            # TextFrameオブジェクトはデフォルトで1つ段落を持つ
    p = tf.add_paragraph()		                           # paragraphオブジェクトの追加作成(2段落目)
    p.text = "This is a second paragraph that's bold"    # textプロパティによる文字列の設定

    # list_to_str = ','.join(conversationHistory)
    list_to_str = ','.join([str(i) for i in conversationHistory])
    # p.text += "\n" + 
    p.text = list_to_str # chatGPT_responce # ,'.join(conversationHistory)  # 区切り文字を「・」にして文字列に変換 # args_text_by_llm # LLMが生成したテキスト
    p.font.bold = True		                               # font.boldプロパティによる太文字設定
    # prs.save('sample.pptx')
    prs.save("C:/Users/0107409377/Desktop/data/sample.pptx")

    import win32com.client as win32
    triggar = True
    def send_mail_based_on_condition(condition):
        try:
            if condition:
                outlook = win32.Dispatch("Outlook.Application")
                mail = outlook.CreateItem(0)
                mail.Subject = "2024/04/26 条件を満たしたため、メールを送信します"
                mail.Body = "これは条件に基づいて送信されたメールです。\nThis report was generated and sent by chatGPT."
                # mail.Body += "\nReport Powered by ChatGPT"
                # mail.To = "test@example.com"
                mail.to = 'Kenji.B.Horiuchi@sony.com; hrucknj@icloud.com'
                mail.cc = 'stmuymte@gmail.com'
                # mail.bcc = 'momoko@mahodo.com'

                # 添付ファイル
                # attachment = "C:/Users/0107409377/Desktop/data/test_mail.txt" # "C:\\path\\to\\your\\file.txt"
                attachment = "C:/Users/0107409377/Desktop/data/sample.pptx" # "C:\\path\\to\\your\\file.txt"
                mail.Attachments.Add(attachment)
                mail.Send()

                print(f"メール送信に成功しました") # : {str(e)}")
        except Exception as e:
            print(f"メール送信に失敗しました: {str(e)}")

    send_mail_based_on_condition(triggar) # True)

            

