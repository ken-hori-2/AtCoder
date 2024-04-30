# from pptx import Presentation
# from pptx.util import Inches

# # テキストファイルからスライドデータを読み込む
# with open("C:/Users/0107409377/Desktop/data/slide-data.txt", "r", encoding="utf-8") as f:
#     slides_data = f.read().splitlines()

# # プレゼンテーションの作成とスライドの生成
# prs = Presentation()

# # リストの長さを2で割った回数までループ
# for i in range(0, len(slides_data)//2):
#     slide = prs.slides.add_slide(prs.slide_layouts[1])
    
#     # リストからタイトルとコンテンツを取得
#     title = slides_data[i * 2]
#     content = slides_data[i * 2 + 1]

#     title_box = slide.shapes.title
#     title_box.text = title

#     content_box = slide.placeholders[1]
#     content_box.text = content

# # ファイル保存
# prs.save("C:/Users/0107409377/Desktop/data/メモ.pptx")


import sys
args_text_by_llm = sys.argv[1]
print("実行時の引数(LLMが生成したテキスト)：", args_text_by_llm)

from pptx import Presentation
prs=Presentation()
from pptx.chart.data import CategoryChartData
c_data=CategoryChartData()
c_data.categories=['A社','B社','C社','D社','E社','F社','G社','H社','I社', 'その他']
c_data.add_series('xxxx（円）',(10,15,8, 5, 17, 12, 2, 8, 15, 7))
sld0=prs.slides.add_slide(prs.slide_layouts[5]) # title only
sld0.shapes[0].text='xxxxに関する市場調査' # 見出し
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

sld0.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    # Cm(1),Cm(3),Cm(20),Cm(15),c_data)
    Cm(1),Cm(3),Cm(15),Cm(15),c_data)

txBox = sld0.shapes.add_textbox(Cm(15),Cm(5),Cm(3),Cm(3)) # left, top, width, height)    #Text Box Shapeオブジェクトの追加
# from pptx.enum.shapes import MSO_SHAPE
# # Slidesコレクション取得
# slides = prs.slides
# # 1枚目のスライドを取得
# slide = slides[0]
# # 末尾にスライドを追加 返り値は追加されたスライド
# # added_slide = slides.add_slide(side_layouts)
# shapes = slide.shapes
# shape = shapes[0]
# for shape in shapes:
#   if shape.shape_type == MSO_SHAPE.RECTANGLE:
#       print("コレは四角形")
# left, top, width, height = shape.left, shape.top, shape.width, shape.height
# text = shape.text
# # # add〇〇系のメソッドで使う
# # shapes.add_shape(autoshape_type_id, left, top, width, height)
# txBox = sld0.shapes.add_textbox(left, top, width, height)

tf = txBox.text_frame		# TextFrameオブジェクトの設定
tf.text = "This is text inside a textbox"            # TextFrameオブジェクトはデフォルトで1つ段落を持つ
p = tf.add_paragraph()		                           # paragraphオブジェクトの追加作成(2段落目)
p.text = "This is a second paragraph that's bold"    # textプロパティによる文字列の設定

p.text += "\n" + args_text_by_llm # LLMが生成したテキスト
p.font.bold = True		                               # font.boldプロパティによる太文字設定

# prs.save('sample.pptx')
prs.save("C:/Users/0107409377/Desktop/data/sample.pptx")

# 図の挿入
# pict_slide_layout=prs.slide_layouts[8]
# sid=prs.slides.add_slide(pict_slide_layout)
# title=sid.placeholders[0]
# pict=sid.placeholders[1]
# body=sid.placeholders[2]

# title.text='サンプル'
# pict.insert_picture('test.png')
# body.text='第2回Pythonでパワポ'
# prs.save("C:/Users/0107409377/Desktop/data/図ありパワポ.pptx")